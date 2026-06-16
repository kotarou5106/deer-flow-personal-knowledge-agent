from __future__ import annotations

import html
import mimetypes
import tempfile
import zipfile
from abc import ABC, abstractmethod
from io import BytesIO
from pathlib import Path
from xml.etree import ElementTree

from bs4 import BeautifulSoup
from openpyxl import load_workbook
from pptx import Presentation

from deerflow.knowledge.ingestion.models import ParsedDocument, TextBlock


class ParserError(RuntimeError):
    pass


class DocumentParser(ABC):
    name: str
    version: str = "1"
    media_types: frozenset[str] = frozenset()
    extensions: frozenset[str] = frozenset()

    @abstractmethod
    def parse(self, data: bytes, *, filename: str | None = None, media_type: str | None = None) -> ParsedDocument:
        raise NotImplementedError


def _offset_blocks(blocks: list[TextBlock]) -> tuple[TextBlock, ...]:
    offset = 0
    result: list[TextBlock] = []
    for block in blocks:
        text = block.text.strip()
        if not text:
            continue
        start = offset
        end = start + len(text)
        result.append(
            TextBlock(
                text=text,
                section_path=block.section_path,
                page_number=block.page_number,
                slide_number=block.slide_number,
                sheet_name=block.sheet_name,
                row_start=block.row_start,
                row_end=block.row_end,
                start_offset=start,
                end_offset=end,
            )
        )
        offset = end + 2
    return tuple(result)


def _decode_text(data: bytes) -> str:
    for encoding in ("utf-8-sig", "utf-8", "utf-16"):
        try:
            return data.decode(encoding)
        except UnicodeDecodeError:
            continue
    raise ParserError("Text file encoding is not supported")


class PlainTextParser(DocumentParser):
    name = "plain_text"
    media_types = frozenset({"text/plain"})
    extensions = frozenset({".txt"})

    def parse(self, data: bytes, *, filename: str | None = None, media_type: str | None = None) -> ParsedDocument:
        text = _decode_text(data).strip()
        if not text:
            raise ParserError("Text file has no extractable content")
        return ParsedDocument(title=filename, text_blocks=_offset_blocks([TextBlock(text=text)]), parser_name=self.name, parser_version=self.version)


class MarkdownParser(DocumentParser):
    name = "markdown"
    media_types = frozenset({"text/markdown", "text/x-markdown"})
    extensions = frozenset({".md", ".markdown"})

    def parse(self, data: bytes, *, filename: str | None = None, media_type: str | None = None) -> ParsedDocument:
        text = _decode_text(data).strip()
        if not text:
            raise ParserError("Markdown file has no extractable content")
        blocks: list[TextBlock] = []
        section: tuple[str, ...] = ()
        current: list[str] = []
        for line in text.splitlines():
            stripped = line.strip()
            if stripped.startswith("#"):
                if current:
                    blocks.append(TextBlock(text="\n".join(current), section_path=section))
                    current = []
                title = stripped.lstrip("#").strip()
                if title:
                    section = (title,)
                    blocks.append(TextBlock(text=title, section_path=section))
                continue
            if stripped:
                current.append(line)
            elif current:
                blocks.append(TextBlock(text="\n".join(current), section_path=section))
                current = []
        if current:
            blocks.append(TextBlock(text="\n".join(current), section_path=section))
        return ParsedDocument(title=filename, text_blocks=_offset_blocks(blocks), parser_name=self.name, parser_version=self.version)


class HtmlParser(DocumentParser):
    name = "html"
    media_types = frozenset({"text/html", "application/xhtml+xml"})
    extensions = frozenset({".html", ".htm", ".xhtml"})

    def parse(self, data: bytes, *, filename: str | None = None, media_type: str | None = None) -> ParsedDocument:
        soup = BeautifulSoup(_decode_text(data), "html.parser")
        for node in soup(["script", "style", "noscript", "template", "svg", "nav", "header", "footer"]):
            node.decompose()
        title = soup.title.get_text(" ", strip=True) if soup.title else filename
        blocks: list[TextBlock] = []
        section: tuple[str, ...] = ()
        for element in soup.find_all(["h1", "h2", "h3", "h4", "p", "li", "td", "th"]):
            text = element.get_text(" ", strip=True)
            if not text:
                continue
            if element.name in {"h1", "h2", "h3", "h4"}:
                section = (text,)
            blocks.append(TextBlock(text=html.unescape(text), section_path=section))
        if not blocks:
            raise ParserError("HTML file has no extractable content")
        return ParsedDocument(title=title, text_blocks=_offset_blocks(blocks), parser_name=self.name, parser_version=self.version)


class DocxParser(DocumentParser):
    name = "docx"
    media_types = frozenset({"application/vnd.openxmlformats-officedocument.wordprocessingml.document"})
    extensions = frozenset({".docx"})

    def parse(self, data: bytes, *, filename: str | None = None, media_type: str | None = None) -> ParsedDocument:
        try:
            with zipfile.ZipFile(BytesIO(data)) as archive:
                xml = archive.read("word/document.xml")
        except Exception as exc:
            raise ParserError("DOCX file is encrypted, corrupt, or unsupported") from exc
        root = ElementTree.fromstring(xml)
        ns = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}
        blocks: list[TextBlock] = []
        section: tuple[str, ...] = ()
        for para in root.findall(".//w:p", ns):
            text = "".join(node.text or "" for node in para.findall(".//w:t", ns)).strip()
            if not text:
                continue
            style = para.find(".//w:pStyle", ns)
            style_value = style.attrib.get(f"{{{ns['w']}}}val", "") if style is not None else ""
            if style_value.lower().startswith("heading"):
                section = (text,)
            blocks.append(TextBlock(text=text, section_path=section))
        if not blocks:
            raise ParserError("DOCX file has no extractable content")
        return ParsedDocument(title=filename, text_blocks=_offset_blocks(blocks), parser_name=self.name, parser_version=self.version)


class PptxParser(DocumentParser):
    name = "pptx"
    media_types = frozenset({"application/vnd.openxmlformats-officedocument.presentationml.presentation"})
    extensions = frozenset({".pptx"})

    def parse(self, data: bytes, *, filename: str | None = None, media_type: str | None = None) -> ParsedDocument:
        try:
            presentation = Presentation(BytesIO(data))
        except Exception as exc:
            raise ParserError("PPTX file is encrypted, corrupt, or unsupported") from exc
        blocks: list[TextBlock] = []
        for slide_number, slide in enumerate(presentation.slides, 1):
            texts = [shape.text.strip() for shape in slide.shapes if getattr(shape, "has_text_frame", False) and shape.text.strip()]
            if not texts:
                continue
            title = texts[0]
            blocks.append(TextBlock(text="\n".join(texts), section_path=(title,), slide_number=slide_number))
        if not blocks:
            raise ParserError("PPTX file has no extractable content")
        return ParsedDocument(title=filename, text_blocks=_offset_blocks(blocks), parser_name=self.name, parser_version=self.version)


class XlsxParser(DocumentParser):
    name = "xlsx"
    media_types = frozenset({"application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"})
    extensions = frozenset({".xlsx"})

    def parse(self, data: bytes, *, filename: str | None = None, media_type: str | None = None) -> ParsedDocument:
        try:
            workbook = load_workbook(BytesIO(data), read_only=True, data_only=True)
        except Exception as exc:
            raise ParserError("XLSX file is encrypted, corrupt, or unsupported") from exc
        blocks: list[TextBlock] = []
        for sheet in workbook.worksheets:
            rows: list[str] = []
            row_start: int | None = None
            row_end: int | None = None
            for row_number, row in enumerate(sheet.iter_rows(values_only=True), 1):
                values = [str(value) for value in row if value is not None and str(value).strip()]
                if not values:
                    continue
                if row_start is None:
                    row_start = row_number
                row_end = row_number
                rows.append("\t".join(values))
            if rows:
                blocks.append(TextBlock(text="\n".join(rows), section_path=(sheet.title,), sheet_name=sheet.title, row_start=row_start, row_end=row_end))
        if not blocks:
            raise ParserError("XLSX file has no extractable content")
        return ParsedDocument(title=filename, text_blocks=_offset_blocks(blocks), parser_name=self.name, parser_version=self.version)


class PdfParser(DocumentParser):
    name = "pdf"
    media_types = frozenset({"application/pdf"})
    extensions = frozenset({".pdf"})

    def parse(self, data: bytes, *, filename: str | None = None, media_type: str | None = None) -> ParsedDocument:
        import pdfplumber

        blocks: list[TextBlock] = []
        with tempfile.NamedTemporaryFile(suffix=".pdf") as tmp:
            tmp.write(data)
            tmp.flush()
            try:
                with pdfplumber.open(tmp.name) as pdf:
                    for page_number, page in enumerate(pdf.pages, 1):
                        text = (page.extract_text() or "").strip()
                        if text:
                            blocks.append(TextBlock(text=text, page_number=page_number, section_path=(f"Page {page_number}",)))
            except Exception as exc:
                raise ParserError("PDF file is encrypted, corrupt, or unsupported") from exc
        if not blocks:
            raise ParserError("PDF file has no extractable content")
        return ParsedDocument(title=filename, text_blocks=_offset_blocks(blocks), parser_name=self.name, parser_version=self.version)


class ParserRegistry:
    def __init__(self, parsers: list[DocumentParser] | None = None) -> None:
        self.parsers = parsers or [PlainTextParser(), MarkdownParser(), HtmlParser(), DocxParser(), PptxParser(), XlsxParser(), PdfParser()]

    def parser_for(self, *, media_type: str | None, filename: str | None, data: bytes) -> DocumentParser:
        normalized_media = (media_type or "").split(";", 1)[0].lower()
        suffix = Path(filename or "").suffix.lower()
        guessed_media = mimetypes.guess_type(filename or "")[0] or ""
        for parser in self.parsers:
            if normalized_media and normalized_media in parser.media_types:
                return parser
        for parser in self.parsers:
            if guessed_media and guessed_media in parser.media_types:
                return parser
        for parser in self.parsers:
            if suffix and suffix in parser.extensions:
                return parser
        if data.startswith(b"%PDF"):
            return next(parser for parser in self.parsers if parser.name == "pdf")
        if data.startswith(b"PK"):
            raise ParserError("Office document type could not be determined")
        raise ParserError("No parser registered for source")

    def parse(self, data: bytes, *, filename: str | None = None, media_type: str | None = None) -> ParsedDocument:
        parser = self.parser_for(media_type=media_type, filename=filename, data=data)
        return parser.parse(data, filename=filename, media_type=media_type)
