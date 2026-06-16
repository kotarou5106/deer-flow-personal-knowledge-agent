from __future__ import annotations

from io import BytesIO
from zipfile import ZipFile

import pytest
from openpyxl import Workbook
from pptx import Presentation

from deerflow.knowledge.ingestion.chunker import ParentChildChunker
from deerflow.knowledge.ingestion.models import ChunkingConfig, ParsedDocument, TextBlock
from deerflow.knowledge.ingestion.parser_registry import ParserError, ParserRegistry


def _docx_bytes() -> bytes:
    document_xml = """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
    <w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">
      <w:body>
        <w:p><w:pPr><w:pStyle w:val="Heading1"/></w:pPr><w:r><w:t>Doc Title</w:t></w:r></w:p>
        <w:p><w:r><w:t>Doc paragraph</w:t></w:r></w:p>
      </w:body>
    </w:document>"""
    buffer = BytesIO()
    with ZipFile(buffer, "w") as archive:
        archive.writestr("[Content_Types].xml", "<Types/>")
        archive.writestr("word/document.xml", document_xml)
    return buffer.getvalue()


def _pptx_bytes() -> bytes:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Slide Title"
    box = slide.shapes.add_textbox(0, 0, 100, 100)
    box.text = "Slide body"
    buffer = BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def _xlsx_bytes() -> bytes:
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Sheet A"
    sheet.append(["Name", "Value"])
    sheet.append(["Alpha", 42])
    empty = workbook.create_sheet("Empty")
    empty.append([])
    buffer = BytesIO()
    workbook.save(buffer)
    return buffer.getvalue()


def _minimal_pdf_bytes() -> bytes:
    objects = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R 5 0 R 6 0 R] /Count 3 >>",
        b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 300 200] /Resources << /Font << /F1 8 0 R >> >> /Contents 4 0 R >>",
        b"<< /Length 44 >>\nstream\nBT /F1 18 Tf 40 120 Td (First PDF Page) Tj ET\nendstream",
        b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 300 200] /Resources << /Font << /F1 8 0 R >> >> >>",
        b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 300 200] /Resources << /Font << /F1 8 0 R >> >> /Contents 7 0 R >>",
        b"<< /Length 45 >>\nstream\nBT /F1 18 Tf 40 120 Td (Second PDF Page) Tj ET\nendstream",
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
    ]
    output = bytearray(b"%PDF-1.4\n%\xe2\xe3\xcf\xd3\n")
    offsets = [0]
    for index, obj in enumerate(objects, 1):
        offsets.append(len(output))
        output.extend(f"{index} 0 obj\n".encode())
        output.extend(obj)
        output.extend(b"\nendobj\n")
    xref = len(output)
    output.extend(f"xref\n0 {len(objects) + 1}\n0000000000 65535 f \n".encode())
    for offset in offsets[1:]:
        output.extend(f"{offset:010d} 00000 n \n".encode())
    output.extend(f"trailer\n<< /Size {len(objects) + 1} /Root 1 0 R >>\nstartxref\n{xref}\n%%EOF\n".encode())
    return bytes(output)


def test_registry_parses_txt_markdown_html_docx_pptx_xlsx() -> None:
    registry = ParserRegistry()

    txt = registry.parse(b"plain text", filename="note.txt")
    md = registry.parse(b"# Title\n\nBody", filename="note.md")
    html = registry.parse(b"<html><head><title>T</title><script>bad()</script></head><body><h1>Head</h1><p>Body</p></body></html>", filename="page.html")
    docx = registry.parse(_docx_bytes(), filename="doc.docx")
    pptx = registry.parse(_pptx_bytes(), filename="deck.pptx")
    xlsx = registry.parse(_xlsx_bytes(), filename="book.xlsx")

    assert txt.text == "plain text"
    assert md.text_blocks[0].section_path == ("Title",)
    assert "bad()" not in html.text
    assert docx.text_blocks[1].section_path == ("Doc Title",)
    assert pptx.text_blocks[0].slide_number == 1
    assert xlsx.text_blocks[0].sheet_name == "Sheet A"
    assert xlsx.text_blocks[0].row_start == 1
    assert xlsx.text_blocks[0].row_end == 2


def test_registry_selects_pdf_parser_by_signature() -> None:
    parser = ParserRegistry().parser_for(media_type=None, filename="unknown.bin", data=b"%PDF-1.4\n")

    assert parser.name == "pdf"


def test_pdf_parser_extracts_text_pages_and_skips_empty_pages() -> None:
    parsed = ParserRegistry().parse(_minimal_pdf_bytes(), filename="sample.pdf")

    assert [block.page_number for block in parsed.text_blocks] == [1, 3]
    assert "First PDF Page" in parsed.text_blocks[0].text
    assert "Second PDF Page" in parsed.text_blocks[1].text


def test_registry_reports_corrupt_or_unsupported_file() -> None:
    with pytest.raises(ParserError):
        ParserRegistry().parse(b"PK not an office document", filename="unknown.zip")

    with pytest.raises(ParserError, match="PDF file is encrypted, corrupt, or unsupported"):
        ParserRegistry().parse(b"%PDF corrupt", filename="broken.pdf")


def test_parent_child_chunking_is_deterministic_and_links_children_to_parent() -> None:
    document = ParsedDocument(
        title="doc",
        parser_name="test",
        parser_version="1",
        text_blocks=(
            TextBlock(text="Section one has enough content to split into children.", section_path=("One",), page_number=1, start_offset=0, end_offset=56),
            TextBlock(text="Section two stays separate.", section_path=("Two",), page_number=2, start_offset=58, end_offset=83),
        ),
    )
    chunker = ParentChildChunker(ChunkingConfig(parent_max_chars=200, child_max_chars=24, child_overlap_chars=4))

    first = chunker.chunk(document)
    second = chunker.chunk(document)

    assert first == second
    assert first[0].parent_index is None
    assert first[1].parent_index == 0
    assert first[-1].section_path == ("Two",)
    assert first[-1].page_number == 2
